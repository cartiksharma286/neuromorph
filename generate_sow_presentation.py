from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


OUT_PATH = Path('/Users/cartiksharma/Downloads/SOW_design_execution.ppt')


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets, level0_font=24, body_font=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    first = True
    for bullet in bullets:
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()

        p.text = text
        p.level = level
        p.font.size = Pt(level0_font if level == 0 else body_font)


def build_presentation():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        'SOW Strategy, Architecture, and Execution Plan',
        'Teqfocus Contractor SOW Draft | Prepared for Lead Architect Delivery | ' + date.today().isoformat(),
    )

    # Slide 2: SOW Scope Summary
    add_bullets_slide(
        prs,
        'SOW Scope and Expected Outcomes',
        [
            'Engagement Focus: Lead Architect support for clinical technology programme (Project Delta).',
            'D1a: Harden AWS target architecture, close gaps, deliver marked-up architecture diagrams and TCO/capability rationale.',
            'D1b: Define Cerebrum Sync App pattern; build development PoC for selected sync approach.',
            'D1c: Stand up FHIR API Gateway on Amazon HealthLake and integrate one external partner end-to-end.',
            'D4: Deliver full phased scope and roadmap with dependencies, assumptions, and decision gates.',
            'D5: Deliver board-ready deck on strategy, risks, investments, and 12-month roadmap.',
        ],
    )

    # Slide 3: Strategy to Approach the SOW
    add_bullets_slide(
        prs,
        'Strategy to Approach the SOW',
        [
            '1) Establish Current-State Baseline and Gaps (Week 1).',
            ('Rapid architecture discovery workshops with Product Owner, operations, and partner teams.', 1),
            ('Document current integrations, latency/SLA constraints, and data governance controls.', 1),
            '2) Finalize Target-State Architecture and Decision Framework (Week 1-2).',
            ('Publish architecture decision records (ADRs), capability matrix, and TCO comparison.', 1),
            '3) Execute Two Technical Proof Points (Week 2-3).',
            ('Cerebrum Sync PoC: CDC via DMS vs EventBridge/Kinesis vs dual-write.', 1),
            ('FHIR Gateway PoC on HealthLake with one partner integrated end-to-end.', 1),
            '4) Consolidate Delivery Plan and Board Narrative (Week 4).',
            ('Roadmap, risks, spend profile, and implementation decision gates.', 1),
        ],
    )

    # Slide 4: Proposed Architecture
    add_bullets_slide(
        prs,
        'Proposed Architecture (AWS Clinical Knowledge Platform)',
        [
            'Data Sources: EMR/EHR systems, partner clinical feeds, PACS/DICOM imaging streams.',
            'Ingestion Layer: API Gateway + Lambda for synchronous APIs; DMS + Kinesis for CDC streams.',
            'Event Backbone: EventBridge for domain events; Kinesis for high-throughput ordered pipelines.',
            'Clinical Data Layer: Amazon HealthLake as FHIR-native canonical store (R4/R4B).',
            'Storage and Analytics: S3 data lake zones for raw/curated datasets and audit retention.',
            'Interoperability Services: SMART on FHIR compliant APIs and partner onboarding templates.',
            'Governance and Security: PHIPA-aware controls, encryption, access policy boundaries, lineage.',
        ],
    )

    # Slide 5: Sync Pattern Decision
    add_bullets_slide(
        prs,
        'Cerebrum Sync App: Pattern Selection Framework',
        [
            'Option A: CDC via AWS DMS -> EventBridge/Kinesis -> HealthLake upsert services.',
            ('Pros: low source impact, replayable events, robust audit trail.', 1),
            ('Considerations: schema drift handling and idempotency controls.', 1),
            'Option B: Native event-driven integration from source systems to EventBridge/Kinesis.',
            ('Pros: lower latency and cleaner event contracts.', 1),
            ('Considerations: partner maturity and contract versioning governance.', 1),
            'Option C: Dual-write service pattern for transactional + FHIR payloads.',
            ('Pros: deterministic write path for specific workflows.', 1),
            ('Considerations: coupling risk and rollback complexity.', 1),
            'Recommendation: Start with DMS+Events for rapid de-risking, then evolve by domain capability.',
        ],
    )

    # Slide 6: Execution Plan
    add_bullets_slide(
        prs,
        'Execution Plan and Milestones',
        [
            'Week 1: Stakeholder alignment, current-state assessment, and architecture gap register.',
            'Week 2: Target-state blueprint, ADRs, and TCO/capability recommendations completed.',
            'Week 3: Two build tracks in parallel:',
            ('Track A: Cerebrum Sync PoC and validation metrics.', 1),
            ('Track B: FHIR gateway and first external partner integration in dev.', 1),
            'Week 4: Consolidated scope + phased roadmap + board presentation package finalized.',
            'Outputs: D1a, D1b, D1c, D4, and D5 ready for approval checkpoints.',
        ],
    )

    # Slide 7: Governance, Risks, and Success Metrics
    add_bullets_slide(
        prs,
        'Governance, Risks, and Success Metrics',
        [
            'Governance: Weekly architecture forum, decision log, and issue/risk escalation cadence.',
            'Primary Risks: partner API variability, data mapping complexity, and timeline compression.',
            'Mitigations: canonical FHIR profile templates, CDC replay strategy, and test harness automation.',
            'Quality Gates: interoperability conformance, latency thresholds, resilience and recovery tests.',
            'Success Metrics: partner integration lead-time, sync accuracy, architecture sign-off, roadmap approval.',
        ],
    )

    # Slide 8: Past Examples Based on CartikSharma.pdf
    add_bullets_slide(
        prs,
        'Relevant Past Examples (from CartikSharma.pdf)',
        [
            'Qmorphix: Led AWS-based clinical/imaging architecture, CDC + event-driven patterns, and TCO-driven recommendations.',
            'Intelerad Medical Systems: Implemented FHIR (R4/R4B) integrations linking DICOM/PACS events to EMR/EHR workflows.',
            'Eigen Health Services: Architected medical imaging exchange with FHIR contracts and external partner integration governance.',
            'Boston Scientific: Served as architecture authority across clinical technology workstreams and integration design decisions.',
            'Neuromorph: Produced phased technology roadmaps and executive strategy materials for cloud clinical platforms.',
            'Healthcare Depth: 8+ years across hospitals/health care and regulated medical technology environments.',
        ],
    )

    # Slide 9: Deliverable Mapping
    add_bullets_slide(
        prs,
        'Deliverable Mapping to SOW Acceptance',
        [
            'D1a Architecture: marked-up diagrams, gap closure register, and AWS TCO/capability recommendation pack.',
            'D1b Sync App: selected synchronization architecture + development PoC + runbook.',
            'D1c FHIR Gateway: HealthLake-based integration module + one external partner pattern proof.',
            'D4 Scope and Roadmap: phased execution plan with dependencies, assumptions, and decision gates.',
            'D5 Board Deck: strategy, risk landscape, investment rationale, and 12-month roadmap narrative.',
        ],
    )

    # Slide 10: Immediate Next Steps
    add_bullets_slide(
        prs,
        'Immediate Next Steps (First 10 Business Days)',
        [
            'Confirm architecture governance model, stakeholder RACI, and sign-off workflow.',
            'Finalize source system inventory and data contract baseline for first partner.',
            'Approve sync pattern evaluation matrix and PoC success criteria.',
            'Launch HealthLake dev environment readiness checks and API security baseline.',
            'Schedule board narrative checkpoint and roadmap review with client stakeholders.',
        ],
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f'Presentation written to {OUT_PATH}')


if __name__ == '__main__':
    build_presentation()
